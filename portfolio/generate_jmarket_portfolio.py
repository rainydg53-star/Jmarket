from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = "Jmarket_Portfolio.pptx"
FONT = "Malgun Gothic"

NAVY = RGBColor(18, 31, 52)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(5, 150, 105)
RED = RGBColor(220, 38, 38)
GRAY = RGBColor(100, 116, 139)
LIGHT = RGBColor(248, 250, 252)
BORDER = RGBColor(203, 213, 225)
BLACK = RGBColor(15, 23, 42)
WHITE = RGBColor(255, 255, 255)


def set_run(run, size=18, bold=False, color=BLACK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text="", size=18, bold=False, color=BLACK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size, bold, color)
    return box


def add_title(slide, title, subtitle=None, section=None):
    add_textbox(slide, 0.55, 0.35, 8.6, 0.45, title, 24, True, NAVY)
    if subtitle:
        add_textbox(slide, 0.58, 0.82, 8.9, 0.3, subtitle, 10, False, GRAY)
    if section:
        add_badge(slide, 10.95, 0.38, 1.4, 0.34, section, BLUE)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.17), Inches(11.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()


def add_badge(slide, x, y, w, h, text, fill=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run(run, 9, True, WHITE)
    return shape


def add_footer(slide, page):
    add_textbox(slide, 0.55, 6.88, 5, 0.18, "Jmarket Portfolio", 8, False, GRAY)
    add_textbox(slide, 11.45, 6.88, 0.55, 0.18, f"{page:02}", 8, False, GRAY, PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, h, items, size=14, color=BLACK, bullet=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.text = item
        if bullet:
            p._p.get_or_add_pPr().set("marL", "228600")
            p._p.get_or_add_pPr().set("indent", "-171450")
        for run in p.runs:
            set_run(run, size, False, color)
        p.space_after = Pt(5)
    return box


def add_panel(slide, x, y, w, h, title=None, fill=LIGHT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    if title:
        add_textbox(slide, x + 0.22, y + 0.18, w - 0.44, 0.28, title, 13, True, NAVY)
    return shape


def add_placeholder(slide, x, y, w, h, text):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(241, 245, 249)
    shape.line.color.rgb = BORDER
    shape.line.dash_style = 4
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run(run, 14, True, GRAY)
    return shape


def add_table(slide, x, y, w, h, rows, col_widths=None, font_size=10):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    set_run(run, font_size, r == 0, WHITE if r == 0 else BLACK)
    return table_shape


def add_process(slide, y, items, colors=None):
    colors = colors or [BLUE, GREEN, RGBColor(124, 58, 237), RGBColor(14, 165, 233), RGBColor(234, 88, 12)]
    x = 0.7
    box_w = 2.0
    for idx, item in enumerate(items):
        fill = colors[idx % len(colors)]
        add_badge(slide, x, y, box_w, 0.56, item, fill)
        if idx < len(items) - 1:
            add_textbox(slide, x + box_w + 0.05, y + 0.14, 0.28, 0.2, "→", 14, True, GRAY)
        x += box_w + 0.35


def make_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    page = 1

    def slide(title, subtitle=None, section=None):
        nonlocal page
        s = prs.slides.add_slide(blank)
        add_title(s, title, subtitle, section)
        add_footer(s, page)
        page += 1
        return s

    # 1
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    add_textbox(s, 0.9, 1.25, 9.8, 0.7, "Jmarket", 36, True, WHITE)
    add_textbox(s, 0.95, 2.0, 10.5, 0.55, "중고거래 및 경매 기반 마켓 플랫폼", 22, True, WHITE)
    add_textbox(s, 0.98, 2.72, 9.8, 0.45, "사용자 거래 흐름과 운영 관리를 함께 설계하는 백엔드/풀스택 개발자", 15, False, RGBColor(226, 232, 240))
    add_bullets(s, 0.98, 3.55, 6.8, 1.4, ["Spring Boot", "React", "WebSocket/STOMP", "마일리지 정산", "관리자 운영"], 13, WHITE)
    add_textbox(s, 0.95, 6.35, 5.5, 0.35, "이름 / 연락처 / 이메일", 12, False, RGBColor(203, 213, 225))
    add_footer(s, page)
    page += 1

    # 2
    s = slide("Profile", "짧고 명확한 인적사항과 개발자 소개", "Profile")
    add_panel(s, 0.7, 1.55, 5.7, 4.65, "인적사항")
    add_table(s, 1.0, 2.1, 5.1, 2.05, [
        ["항목", "내용"],
        ["이름", ""],
        ["연락처", ""],
        ["이메일", ""],
        ["GitHub", ""],
        ["Blog / Notion", ""],
    ], [1.35, 3.75], 10)
    add_panel(s, 6.75, 1.55, 5.6, 4.65, "소개")
    add_bullets(s, 7.05, 2.15, 4.95, 2.4, [
        "Spring Boot와 React 기반 개인 프로젝트 개발",
        "중고거래, 경매, 채팅, 마일리지, 관리자 운영 기능 구현",
        "기능 나열보다 권한, 상태 전이, 데이터 정합성을 고려한 설계 경험",
    ], 13)

    # 3
    s = slide("Education / Skills", "기술 스택은 사용 이유가 드러나도록 정리", "Profile")
    add_table(s, 0.75, 1.45, 5.5, 1.45, [["기간", "기관", "내용"], ["", "", ""], ["", "", ""]], [1.35, 1.8, 2.35], 9)
    add_table(s, 6.6, 1.45, 5.65, 3.4, [
        ["구분", "기술"],
        ["Backend", "Java 21, Spring Boot, Spring Security, JPA"],
        ["Frontend", "React, Vite, JavaScript, CSS"],
        ["DB / Cache", "MySQL, Redis"],
        ["Realtime", "WebSocket, STOMP"],
        ["Test / Tool", "JUnit, MockMvc, Gradle, npm, Git"],
    ], [1.35, 4.3], 8)
    add_table(s, 0.75, 3.35, 5.5, 1.2, [["취득일", "자격증", "기관"], ["", "", ""], ["", "", ""]], [1.35, 2.3, 1.85], 9)
    add_panel(s, 0.75, 5.05, 11.5, 1.25, "개인 학습 / 활동")
    add_bullets(s, 1.0, 5.55, 10.8, 0.55, [
        "인증/인가, 경매 정산, React Hook 의존성, 관리자 운영 화면 중심으로 학습 및 구현",
    ], 12)

    # 4
    s = slide("Project List", "포트폴리오에 포함할 프로젝트 목차", "Contents")
    add_table(s, 0.7, 1.6, 11.9, 1.4, [
        ["프로젝트명", "기간", "인원", "주요 기술", "설명"],
        ["Jmarket", "2026.04 ~ 2026.05", "개인", "Spring Boot, React, MySQL, Redis, WebSocket", "중고거래와 경매를 함께 제공하는 마켓 플랫폼"],
    ], [1.55, 1.65, 0.9, 3.55, 4.25], 9)
    add_placeholder(s, 0.9, 3.45, 11.2, 2.2, "프로젝트 대표 화면 또는 로고 캡처 삽입")

    # 5
    s = slide("Project Overview", "OO 문제를 해결하기 위해 OO 시스템 개발", "Jmarket")
    add_panel(s, 0.7, 1.45, 12.0, 1.05, "핵심 문장")
    add_textbox(s, 1.0, 1.92, 11.2, 0.35, "중고거래 과정에서 발생하는 상품 탐색, 거래 요청, 경매 입찰, 채팅, 마일리지 정산, 관리자 운영 문제를 해결하기 위해 통합 마켓 플랫폼을 개발했습니다.", 13, True, NAVY)
    add_panel(s, 0.7, 2.75, 5.7, 3.55, "개발 배경")
    add_bullets(s, 1.0, 3.28, 5.05, 1.85, [
        "상품 거래, 경매, 채팅, 신고, 정산, 관리자 운영 흐름을 하나의 서비스로 연결",
        "일반 사용자와 관리자 관점을 모두 고려한 거래 플랫폼 설계",
    ], 12)
    add_panel(s, 6.75, 2.75, 5.95, 3.55, "해결하려는 문제")
    add_bullets(s, 7.05, 3.28, 5.35, 2.1, [
        "거래 요청부터 완료까지의 상태 관리",
        "경매 입찰과 마일리지 예약 정합성 유지",
        "신고/출금/회원 제재 등 관리자 운영 분리",
        "권한별 메뉴와 API 접근 제어",
    ], 12)

    # 6
    s = slide("Main Features", "핵심 기능과 기대 효과", "Jmarket")
    add_panel(s, 0.7, 1.45, 5.9, 4.9, "주요 기능")
    add_bullets(s, 1.0, 1.95, 5.3, 3.75, [
        "회원가입, 로그인, 이메일 인증, 비밀번호 찾기/변경",
        "상품 등록/수정/삭제, 검색/필터, 이미지 드래그 업로드",
        "경매 등록, 입찰, 즉시구매, 입찰 내역, 마감 처리",
        "거래 요청/수락/취소/완료, 실시간 채팅",
        "마일리지 충전/사용/출금, 신고/문의/리뷰/알림",
        "관리자 대시보드 및 운영 관리",
    ], 11)
    add_panel(s, 6.95, 1.45, 5.35, 2.25, "기대 효과")
    add_bullets(s, 7.25, 1.95, 4.8, 1.25, [
        "상품 탐색부터 거래 완료까지 한 흐름에서 처리",
        "관리자는 운영 데이터를 빠르게 확인",
        "권한/상태별 액션 노출로 오류 가능성 감소",
    ], 11)
    add_panel(s, 6.95, 4.1, 5.35, 2.25, "차별화 포인트")
    add_bullets(s, 7.25, 4.6, 4.8, 1.25, [
        "일반 상품 거래와 경매 거래를 함께 제공",
        "마일리지 예약/해제/정산을 포함한 입찰 흐름",
        "사용자 화면과 관리자 화면 UX 목적 분리",
    ], 11)

    # 7
    s = slide("Role / Tech Stack", "개인 프로젝트 담당 범위와 개발 환경", "Jmarket")
    add_panel(s, 0.7, 1.45, 6.0, 4.9, "담당 역할")
    add_bullets(s, 1.0, 1.95, 5.35, 3.9, [
        "기획, DB 설계, 백엔드 API, 프론트 화면 구현 전담",
        "Spring Security 기반 인증/인가 구현",
        "상품/경매/거래/채팅/마일리지 도메인 설계",
        "WebSocket/STOMP 채팅 구현",
        "관리자 대시보드 및 액션 모달/토스트 통일",
        "이미지 업로드 서버 검증 및 통합 테스트 작성",
        "React Hook dependency warning 17개 제거",
    ], 11)
    add_table(s, 7.05, 1.55, 5.1, 3.9, [
        ["구분", "내용"],
        ["Backend", "Spring Boot, Java 21, Security, JPA"],
        ["Frontend", "React, Vite, JavaScript"],
        ["DB / Cache", "MySQL, Redis"],
        ["Realtime", "WebSocket, STOMP"],
        ["Build / Test", "Gradle, npm, JUnit, MockMvc"],
    ], [1.4, 3.7], 8)

    # 8-12 Screens
    screen_slides = [
        ("Main Page", "비로그인 사용자도 접근 가능한 메인 화면", "메인 화면 캡처 삽입", ["실시간 검색어 순위", "급상승 물품", "빈 상태 UI", "상품/경매 탐색 흐름"]),
        ("Product Flow", "상품 목록/상세/등록 흐름", "상품 목록 / 상세 / 등록 화면 캡처 삽입", ["검색/필터 조건 URL 유지", "썸네일 클릭 상세 이동", "이미지 드래그 업로드", "서버 이미지 검증"]),
        ("Auction Flow", "경매 목록/상세/입찰 흐름", "경매 목록 / 상세 / 입찰 내역 화면 캡처 삽입", ["상태별 입찰 가능 여부 제어", "최고 입찰자와 최고가 표시", "입찰 내역 차트", "마감 타이머 상태 갱신"]),
        ("Chat / Trade Flow", "채팅과 거래 상태 흐름", "채팅 목록 / 팝업 / 거래 목록 화면 캡처 삽입", ["WebSocket/STOMP 실시간 메시지", "채팅방 새 창 팝업", "거래 요청/수락/취소/완료", "거래 완료 후 리뷰 작성"]),
        ("Admin Dashboard", "운영 도구형 관리자 화면", "관리자 대시보드 / 상품 / 경매 / 출금 화면 캡처 삽입", ["서버 관리자 권한 검증", "검색/필터/테이블/상태 배지", "액션 결과 모달/토스트", "위험 액션 확인 모달"]),
    ]
    for title, sub, ph, bullets in screen_slides:
        s = slide(title, sub, "Screens")
        add_placeholder(s, 0.7, 1.45, 7.0, 4.9, ph)
        add_panel(s, 8.05, 1.45, 4.25, 4.9, "핵심 포인트")
        add_bullets(s, 8.35, 2.05, 3.65, 2.4, bullets, 12)

    # 13 Architecture
    s = slide("System Architecture", "Frontend / Backend / DB / Infra 구조", "Architecture")
    add_process(s, 1.55, ["React Frontend", "REST API", "Spring Boot", "MySQL", "Redis"])
    add_panel(s, 0.85, 2.65, 3.0, 2.5, "Frontend")
    add_bullets(s, 1.1, 3.15, 2.35, 1.2, ["사용자 화면", "관리자 화면", "채팅 팝업"], 11)
    add_panel(s, 4.05, 2.65, 3.7, 2.5, "Backend")
    add_bullets(s, 4.3, 3.15, 3.05, 1.45, ["Auth / Product / Auction / Trade", "Chat / Mileage / Report / Admin", "권한 및 상태 검증"], 11)
    add_panel(s, 8.0, 2.65, 2.0, 2.5, "MySQL")
    add_bullets(s, 8.25, 3.15, 1.45, 1.3, ["회원", "상품/경매", "거래/정산"], 11)
    add_panel(s, 10.25, 2.65, 2.0, 2.5, "Redis")
    add_bullets(s, 10.5, 3.15, 1.45, 1.0, ["이메일 인증", "캐시성 데이터"], 11)
    add_panel(s, 0.85, 5.55, 11.4, 0.75, "시스템 흐름")
    add_textbox(s, 1.1, 5.92, 10.8, 0.25, "React 화면 요청 → Spring Boot 인증/인가 및 상태 검증 → MySQL 저장 → Redis 인증/캐시 → WebSocket 실시간 채팅", 11, True, NAVY)

    # 14 API
    s = slide("Core API", "전체 API가 아닌 핵심 흐름 중심", "API")
    add_table(s, 0.7, 1.45, 11.9, 1.85, [
        ["기능", "Method", "API", "설명"],
        ["이미지 업로드", "POST", "/api/products/images", "상품/경매 등록용 이미지 업로드 및 서버 검증"],
        ["경매 입찰", "POST", "/api/auctions/{auctionId}/bids", "경매 상태, 최고가, 마일리지 검증 후 입찰 처리"],
        ["출금 처리", "PATCH", "/api/admin/mileage/withdrawals/{id}", "관리자 출금 승인/반려 처리"],
    ], [1.45, 1.0, 3.8, 5.65], 8)
    add_panel(s, 0.7, 3.75, 11.9, 1.85, "API 설계 기준")
    add_bullets(s, 1.0, 4.25, 11.0, 0.9, [
        "사용자 액션은 서버에서 권한과 상태를 재검증",
        "관리자 API는 프론트 메뉴 숨김과 별개로 서버에서 권한 강제",
        "금전성 데이터가 포함된 기능은 트랜잭션 경계를 명확히 분리",
    ], 12)

    # 15 DB
    s = slide("DB Design", "주요 도메인과 테이블 구조", "DB")
    add_table(s, 0.7, 1.45, 6.25, 4.4, [
        ["도메인", "주요 테이블", "설명"],
        ["Auth", "users", "회원 정보, 권한, 상태"],
        ["Product", "products, product_images", "상품, 이미지, 찜"],
        ["Auction", "auctions, bids", "경매, 입찰 내역"],
        ["Trade", "trades", "상품 거래 상태"],
        ["Chat", "chat_rooms, chat_messages", "채팅방, 메시지"],
        ["Mileage", "mileage_accounts, ledger, withdrawals", "잔액, 원장, 출금"],
        ["Admin", "reports, restrictions, audit_logs", "신고, 제재, 감사 로그"],
    ], [1.15, 2.75, 2.35], 7)
    add_placeholder(s, 7.25, 1.45, 5.0, 3.1, "ERD 캡처 삽입")
    add_panel(s, 7.25, 4.85, 5.0, 1.0, "설계 포인트")
    add_bullets(s, 7.5, 5.23, 4.45, 0.35, ["상품/경매는 분리하고 사용자, 리뷰, 채팅 흐름은 연결"], 10)

    # 16-19 core logic
    logic_slides = [
        ("Core Logic 1", "경매 입찰 및 마일리지 예약 처리", ["경매 입찰은 현재 최고가, 입찰자 권한, 경매 상태, 마일리지 잔액, 이전 입찰자의 예약 해제를 함께 처리해야 합니다.", "상태/마감/최소 입찰가/본인 상품 입찰/현재 최고 입찰자 여부를 검증합니다.", "신규 최고 입찰자의 마일리지를 예약하고 이전 최고 입찰자의 예약 마일리지를 해제합니다.", "결과적으로 최고가, 최고 입찰자, 예약 금액의 일관성을 유지했습니다."]),
        ("Core Logic 2", "권한 기반 메뉴 및 API 접근 제어", ["프론트에서 버튼을 숨기는 것만으로는 보안이 보장되지 않습니다.", "사용자 역할에 따라 메뉴와 버튼을 제어하고, 관리자 API는 서버에서 접근을 강제했습니다.", "신고/출금/회원 제재 같은 관리자 액션은 서버에서 재검증합니다.", "UI 편의성과 서버 보안을 분리하여 권한 없는 접근을 방지했습니다."]),
        ("Core Logic 3", "이미지 업로드 서버 검증", ["프론트 검증만으로는 확장자를 위장한 파일 업로드를 막기 어렵습니다.", "개수, 용량, 확장자, Content-Type, 파일 시그니처를 서버에서 검증했습니다.", "JPEG/PNG/GIF/WEBP 시그니처를 확인하고 저장 경로도 검증했습니다.", "정상 업로드, 위장 파일 차단, 미지원 확장자 차단 테스트를 작성했습니다."]),
        ("Core Logic 4", "React 데이터 로딩 흐름 안정화", ["Hook dependency warning은 데이터 로딩 구조를 점검하는 신호였습니다.", "조회 함수는 useCallback으로 안정화하고, 검색/필터는 URL query 기반으로 복원했습니다.", "경매 타이머는 auction 객체 전체가 아닌 id/status/endAt만 의존하도록 정리했습니다.", "warning 17개 제거 후 lint/build를 통과했습니다."]),
    ]
    for title, sub, bullets in logic_slides:
        s = slide(title, sub, "Logic")
        add_panel(s, 0.8, 1.55, 11.8, 4.75, sub)
        add_bullets(s, 1.15, 2.15, 11.0, 2.8, bullets, 13)

    # 20-22 troubleshooting
    trouble = [
        ("Troubleshooting 1", "경매 입찰 데이터 정합성", "여러 사용자가 입찰할 때 최고 입찰자, 현재 최고가, 마일리지 예약 상태가 불일치할 가능성이 있었습니다.", "입찰 등록, 최고가 갱신, 마일리지 예약/해제 처리가 서로 다른 도메인에 걸쳐 있었습니다.", "입찰 처리 흐름을 하나의 도메인 흐름으로 묶고, 현재 경매 상태와 최고가를 재검증한 뒤 예약/해제를 처리했습니다.", "동시 입찰 상황에서도 최고 입찰자, 현재 최고가, 예약 금액이 일관되게 유지되도록 개선했습니다."),
        ("Troubleshooting 2", "관리자 권한 서버 강제", "관리자 메뉴와 버튼을 숨기더라도 URL 직접 접근이나 API 직접 호출로 관리자 기능 접근을 시도할 수 있었습니다.", "프론트 UI 제어는 사용자 경험 측면의 처리일 뿐 실제 보안 경계는 서버에서 강제되어야 했습니다.", "Spring Security에서 관리자 API 경로를 제한하고, 프론트에서는 권한별 메뉴/버튼 노출을 분리해 처리했습니다.", "권한 없는 사용자의 직접 접근과 API 호출이 서버에서 차단되도록 개선했습니다."),
        ("Troubleshooting 3", "경매 마감 상태 불일치", "경매 종료 시간이 지난 직후 서버 상태 갱신 전에는 화면상 입찰 가능 상태처럼 보일 수 있었습니다.", "경매 상태는 서버의 OPEN/CLOSED 값과 클라이언트 현재 시간 기준 마감 여부를 함께 고려해야 했습니다.", "프론트에서는 endAt 기준으로 버튼과 상태 배지를 즉시 비활성화하고, 서버 입찰 API에서도 상태와 마감 시간을 재검증했습니다.", "마감 직후에도 화면 상태와 실제 입찰 가능 여부가 일관되도록 개선했습니다."),
    ]
    for title, sub, problem, cause, solution, result in trouble:
        s = slide(title, sub, "Trouble")
        labels = [("문제 상황", problem, RED), ("원인", cause, RGBColor(234, 88, 12)), ("해결 방법", solution, BLUE), ("결과", result, GREEN)]
        y = 1.45
        for label, body, color in labels:
            add_badge(s, 0.8, y, 1.35, 0.38, label, color)
            add_textbox(s, 2.35, y + 0.02, 9.6, 0.45, body, 11, False, BLACK)
            y += 1.1

    # 23 Retrospective
    s = slide("Retrospective", "배운 점, 아쉬운 점, 개선 방향", "Review")
    add_panel(s, 0.7, 1.45, 3.75, 4.9, "배운 점")
    add_bullets(s, 1.0, 1.95, 3.15, 2.6, ["연결된 기능은 상태 전이 설계가 중요", "권한 제어는 서버에서 강제해야 함", "Hook warning은 데이터 로딩 구조 점검 신호"], 11)
    add_panel(s, 4.85, 1.45, 3.75, 4.9, "아쉬운 점")
    add_bullets(s, 5.15, 1.95, 3.15, 2.6, ["전체 E2E 테스트 부족", "Redis 외부 의존성 테스트 환경 분리 필요", "배포 환경 모니터링/로그 설계 보완 필요"], 11)
    add_panel(s, 9.0, 1.45, 3.2, 4.9, "개선 방향")
    add_bullets(s, 9.3, 1.95, 2.6, 2.6, ["경매 동시성 테스트 강화", "테스트 프로파일 분리", "관리자 통계 고도화", "E2E 테스트 추가"], 11)

    # 24 Test
    s = slide("Test / Verification", "구현 후 검증한 내용", "Verify")
    add_table(s, 0.8, 1.55, 11.6, 2.5, [
        ["구분", "검증"],
        ["Frontend", "npm run lint, npm run build"],
        ["Backend", "이미지 업로드 통합 테스트"],
        ["Image Upload", "정상 PNG 업로드, 위장 파일 차단, SVG 차단"],
        ["Hook Cleanup", "기존 17개 warning 제거 후 빌드 성공"],
    ], [2.2, 9.4], 10)
    add_panel(s, 0.8, 4.55, 11.6, 1.2, "남은 보완")
    add_bullets(s, 1.1, 5.05, 11.0, 0.45, ["전체 백엔드 테스트는 Redis 연결 환경 필요", "결제/출금/경매 동시성 관련 테스트 추가 필요"], 11)

    # 25 Flow
    s = slide("Service Flow", "사용자 흐름과 관리자 운영 흐름", "Flow")
    add_panel(s, 0.8, 1.45, 11.7, 1.15, "상품 거래 흐름")
    add_process(s, 1.85, ["상품 등록", "상세 조회", "거래 요청", "채팅", "거래 완료"])
    add_panel(s, 0.8, 3.0, 11.7, 1.15, "경매 거래 흐름")
    add_process(s, 3.4, ["경매 등록", "입찰", "마일리지 예약", "마감", "정산/리뷰"])
    add_panel(s, 0.8, 4.55, 11.7, 1.15, "관리자 운영 흐름")
    add_process(s, 4.95, ["신고/출금", "관리자 검토", "승인/반려", "제재", "로그 관리"])

    # 26 Security
    s = slide("Security Points", "보안과 안정성을 위해 고려한 부분", "Security")
    add_bullets(s, 1.0, 1.65, 11.0, 3.8, [
        "Spring Security 기반 인증/인가",
        "관리자 API 서버 권한 검증",
        "이미지 업로드 파일 형식 검증",
        "사용자 상태에 따른 기능 제한",
        "삭제/취소/차단/반려 등 위험 액션 확인 모달",
        "비밀번호 암호화 저장",
        "이메일 인증 기반 회원가입/비밀번호 찾기",
    ], 15)

    # 27 Future
    s = slide("Future Improvements", "추가 개선 방향", "Future")
    add_bullets(s, 1.0, 1.7, 11.0, 3.5, [
        "배포 환경 구성 및 CI/CD 적용",
        "경매 입찰 동시성 테스트 강화",
        "관리자 통계 시각화 고도화",
        "검색 성능 개선 및 인덱싱 전략 보완",
        "알림 기능 실시간화",
        "주요 사용자 흐름 E2E 테스트 추가",
    ], 15)

    # 28 Closing
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    add_textbox(s, 0.9, 0.9, 8.0, 0.6, "Thank You", 34, True, WHITE)
    add_textbox(s, 0.95, 1.75, 10.8, 0.45, "단순 기능 구현을 넘어 사용자 거래 흐름과 관리자 운영 흐름을 연결한 프로젝트입니다.", 16, False, RGBColor(226, 232, 240))
    add_table(s, 1.0, 3.0, 7.2, 2.0, [
        ["항목", "내용"],
        ["GitHub", ""],
        ["Blog / Notion", ""],
        ["시연 영상", ""],
        ["연락처", ""],
    ], [1.7, 5.5], 11)
    add_footer(s, page)

    return prs


if __name__ == "__main__":
    presentation = make_prs()
    presentation.save(OUT)
    print(OUT)
