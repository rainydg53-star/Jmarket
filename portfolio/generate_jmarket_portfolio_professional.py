from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = "Jmarket_Portfolio_Professional.pptx"
FONT = "Malgun Gothic"

INK = RGBColor(17, 24, 39)
MUTED = RGBColor(100, 116, 139)
LINE = RGBColor(226, 232, 240)
PANEL = RGBColor(248, 250, 252)
NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(5, 150, 105)
ORANGE = RGBColor(234, 88, 12)
RED = RGBColor(220, 38, 38)
WHITE = RGBColor(255, 255, 255)


def rgb(hex_value):
    hex_value = hex_value.lstrip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def fmt(run, size=12, bold=False, color=INK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def text(slide, x, y, w, h, value, size=12, bold=False, color=INK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = value
    fmt(run, size, bold, color)
    return box


def rect(slide, x, y, w, h, fill=WHITE, border=LINE, radius=True):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.7)
    return shape


def title(slide, page, heading, kicker=None):
    text(slide, 0.62, 0.36, 1.35, 0.25, kicker or "Jmarket", 8.5, True, BLUE)
    text(slide, 0.62, 0.62, 8.3, 0.45, heading, 21, True, NAVY)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.18), Inches(12.1), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    text(slide, 11.85, 0.58, 0.75, 0.22, f"{page:02}", 8, False, MUTED, PP_ALIGN.RIGHT)


def section(slide, label, heading, body):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    text(slide, 0.9, 1.35, 2.5, 0.3, label, 10, True, RGBColor(147, 197, 253))
    text(slide, 0.9, 1.85, 9.3, 0.75, heading, 30, True, WHITE)
    text(slide, 0.92, 2.8, 8.6, 0.5, body, 14, False, RGBColor(203, 213, 225))
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(4.0), Inches(1.4), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()


def bullets(slide, x, y, w, h, items, size=11, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(5)
        p._p.get_or_add_pPr().set("marL", "228600")
        p._p.get_or_add_pPr().set("indent", "-171450")
        for r in p.runs:
            fmt(r, size, False, color)
    return box


def numbered(slide, x, y, w, h, items, size=11):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = f"{idx + 1}. {item}"
        p.space_after = Pt(5)
        for r in p.runs:
            fmt(r, size, False, INK)
    return box


def card(slide, x, y, w, h, heading, body=None, accent=BLUE):
    rect(slide, x, y, w, h, WHITE, LINE)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    text(slide, x + 0.22, y + 0.2, w - 0.44, 0.28, heading, 12, True, NAVY)
    if body:
        text(slide, x + 0.22, y + 0.58, w - 0.44, h - 0.75, body, 10.2, False, MUTED)


def placeholder(slide, x, y, w, h, label):
    rect(slide, x, y, w, h, PANEL, LINE)
    text(slide, x + 0.2, y + h / 2 - 0.16, w - 0.4, 0.35, label, 13, True, MUTED, PP_ALIGN.CENTER)


def table(slide, x, y, w, h, rows, widths=None, size=8.5):
    shp = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shp.table
    if widths:
        for idx, width in enumerate(widths):
            tbl.columns[idx].width = Inches(width)
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = value
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r_idx == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    fmt(run, size, r_idx == 0, WHITE if r_idx == 0 else INK)
    return shp


def metric(slide, x, y, w, label, value, color=BLUE):
    rect(slide, x, y, w, 0.95, WHITE, LINE)
    text(slide, x + 0.18, y + 0.18, w - 0.36, 0.2, label, 8.5, True, MUTED)
    text(slide, x + 0.18, y + 0.46, w - 0.36, 0.32, value, 18, True, color)


def flow(slide, x, y, items):
    step_w = 1.8
    for idx, item in enumerate(items):
        rect(slide, x + idx * 2.15, y, step_w, 0.65, WHITE, LINE)
        text(slide, x + idx * 2.15 + 0.12, y + 0.19, step_w - 0.24, 0.22, item, 9.5, True, NAVY, PP_ALIGN.CENTER)
        if idx < len(items) - 1:
            text(slide, x + idx * 2.15 + step_w + 0.07, y + 0.2, 0.25, 0.2, "→", 12, True, MUTED)


def prs_new():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    page = 1

    def add(heading, kicker=None):
        nonlocal page
        s = prs.slides.add_slide(blank)
        title(s, page, heading, kicker)
        page += 1
        return s

    # 1 Cover
    s = prs.slides.add_slide(blank)
    left = slide_bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(4.05), Inches(7.5))
    left.fill.solid(); left.fill.fore_color.rgb = NAVY; left.line.fill.background()
    right = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(4.05), 0, Inches(9.283), Inches(7.5))
    right.fill.solid(); right.fill.fore_color.rgb = WHITE; right.line.fill.background()
    text(s, 0.65, 0.72, 2.8, 0.3, "PORTFOLIO", 10, True, RGBColor(147, 197, 253))
    text(s, 0.65, 1.2, 2.8, 0.7, "Jmarket", 28, True, WHITE)
    text(s, 0.68, 1.92, 2.7, 0.65, "중고거래 및 경매 기반 마켓 플랫폼", 13, False, RGBColor(203, 213, 225))
    text(s, 4.85, 1.25, 7.2, 0.55, "사용자 거래 흐름과 관리자 운영 흐름을 연결한 개인 프로젝트", 24, True, NAVY)
    text(s, 4.88, 2.1, 7.1, 0.55, "Spring Boot · React · WebSocket · MySQL · Redis", 13, False, MUTED)
    card(s, 4.9, 3.1, 2.25, 1.25, "Role", "Full-stack\n개인 프로젝트", BLUE)
    card(s, 7.45, 3.1, 2.25, 1.25, "Focus", "거래 상태\n정합성", GREEN)
    card(s, 10.0, 3.1, 2.25, 1.25, "Domain", "중고거래\n경매/정산", ORANGE)
    text(s, 4.9, 6.5, 5.6, 0.25, "이름 / 연락처 / 이메일", 10.5, False, MUTED)
    page += 1

    # 2-4 Profile/TOC
    s = add("Profile", "01 Profile")
    table(s, 0.75, 1.55, 5.55, 2.25, [["항목", "내용"], ["이름", ""], ["연락처", ""], ["이메일", ""], ["GitHub", ""], ["Blog / Notion", ""]], [1.3, 4.25], 9.5)
    card(s, 6.75, 1.55, 5.7, 2.25, "한 줄 소개", "Spring Boot와 React를 기반으로 거래, 경매, 채팅, 마일리지, 관리자 운영 기능을 구현했습니다. 기능 구현뿐 아니라 권한, 상태 전이, 데이터 정합성을 함께 고려했습니다.", BLUE)
    table(s, 0.75, 4.35, 5.55, 1.15, [["기간", "기관", "내용"], ["", "", ""], ["", "", ""]], [1.3, 1.75, 2.5], 8.5)
    table(s, 6.75, 4.35, 5.7, 1.15, [["취득일", "자격증", "기관"], ["", "", ""], ["", "", ""]], [1.4, 2.4, 1.9], 8.5)

    s = add("Skills", "01 Profile")
    table(s, 0.75, 1.55, 11.8, 3.2, [
        ["구분", "기술", "사용 이유"],
        ["Backend", "Java 21, Spring Boot, Spring Security, JPA", "인증/인가, 트랜잭션, 도메인 API 구현"],
        ["Frontend", "React, Vite, JavaScript", "사용자/관리자 화면과 상태 기반 UI 구현"],
        ["Database / Cache", "MySQL, Redis", "핵심 데이터 저장, 이메일 인증/캐시성 데이터 처리"],
        ["Realtime", "WebSocket, STOMP", "거래/경매 채팅 메시지 실시간 송수신"],
        ["Test / Tool", "JUnit, MockMvc, Gradle, npm, Git", "기능 검증, 빌드, 버전 관리"],
    ], [1.55, 4.1, 6.15], 8)
    card(s, 0.75, 5.25, 11.8, 0.8, "학습/활동 요약", "Spring Security 인증/인가, 경매 정산 흐름, React 데이터 로딩 구조, 관리자 운영 UX를 중심으로 학습하고 프로젝트에 적용했습니다.", GREEN)

    s = add("Project Contents", "02 Contents")
    table(s, 0.75, 1.6, 11.8, 1.3, [["프로젝트명", "기간", "인원", "기술", "설명"], ["Jmarket", "2026.04 ~ 2026.05", "개인", "Spring Boot, React, MySQL, Redis, WebSocket", "중고거래와 경매를 함께 제공하는 마켓 플랫폼"]], [1.5, 1.65, 0.85, 3.65, 4.15], 8.5)
    placeholder(s, 0.75, 3.35, 11.8, 2.35, "프로젝트 대표 화면 캡처 삽입")

    # 5 Section
    s = prs.slides.add_slide(blank); section(s, "03 PROJECT", "Jmarket 상세", "문제 해결 중심으로 프로젝트의 구조와 핵심 로직을 정리했습니다."); page += 1

    # 6 Overview
    s = add("Project Overview", "03 Project")
    card(s, 0.75, 1.5, 11.8, 0.95, "핵심 문장", "중고거래 과정의 상품 탐색, 거래 요청, 경매 입찰, 채팅, 마일리지 정산, 관리자 운영 문제를 해결하기 위해 통합 마켓 플랫폼을 개발했습니다.", BLUE)
    card(s, 0.75, 2.85, 3.75, 2.4, "해결하려는 문제", "거래 요청부터 완료까지 상태 관리\n경매 입찰과 마일리지 예약 정합성\n신고/출금/회원 제재 운영 분리", RED)
    card(s, 4.85, 2.85, 3.75, 2.4, "기대 효과", "사용자는 한 흐름에서 거래 처리\n관리자는 운영 데이터를 빠르게 확인\n권한/상태별 액션으로 오류 감소", GREEN)
    card(s, 8.95, 2.85, 3.6, 2.4, "차별화 포인트", "일반 거래와 경매 동시 제공\n마일리지 예약/해제/정산 포함\n사용자/관리자 UX 분리", ORANGE)

    # 7 Features
    s = add("Main Features", "03 Project")
    features = ["회원가입, 로그인, 이메일 인증, 비밀번호 찾기/변경", "상품 등록/수정/삭제, 이미지 드래그 업로드, 검색/필터", "경매 등록, 입찰, 즉시구매, 입찰 내역, 마감 처리", "거래 요청/수락/취소/완료 및 리뷰", "실시간 채팅방 목록/팝업", "마일리지 충전/사용/출금 요청 및 관리자 처리", "신고, 고객센터 문의, 알림", "관리자 대시보드와 운영 관리"]
    numbered(s, 0.95, 1.55, 11.6, 4.4, features, 12.5)

    # 8 Role
    s = add("Role & Responsibility", "03 Project")
    metric(s, 0.75, 1.45, 2.6, "Project", "개인")
    metric(s, 3.65, 1.45, 2.6, "Slides", "28p")
    metric(s, 6.55, 1.45, 2.6, "Hook Warning", "17 → 0", GREEN)
    metric(s, 9.45, 1.45, 2.6, "Domain", "8+")
    bullets(s, 0.95, 2.85, 11.2, 2.8, ["기획, DB 설계, 백엔드 API, 프론트 화면 구현 전담", "Spring Security 기반 인증/인가와 관리자 API 접근 제한", "상품/경매/거래/채팅/마일리지 도메인 구현", "WebSocket/STOMP 채팅과 관리자 운영 화면 개선", "이미지 업로드 서버 검증 및 통합 테스트 작성"], 12.5)

    # 9-13 Screens
    screens = [
        ("Main Page", "비로그인 사용자도 서비스 콘텐츠를 확인할 수 있는 메인 화면", "메인 화면 캡처 삽입", ["실시간 검색어 순위", "급상승 물품", "빈 상태 UI", "상품/경매 탐색 흐름"]),
        ("Product Flow", "상품 탐색, 등록, 상세 조회 흐름", "상품 목록 / 상세 / 등록 화면 캡처 삽입", ["검색/필터 URL 유지", "썸네일 상세 이동", "드래그 이미지 업로드", "서버 이미지 검증"]),
        ("Auction Flow", "경매 입찰과 마감 상태 확인 흐름", "경매 목록 / 상세 / 입찰 화면 캡처 삽입", ["상태별 입찰 제어", "최고 입찰자/최고가 표시", "입찰 차트", "마감 타이머 갱신"]),
        ("Chat / Trade Flow", "거래 진행 중 대화와 상태 처리", "채팅 목록 / 팝업 / 거래 목록 캡처 삽입", ["STOMP 실시간 메시지", "채팅방 새 창 팝업", "거래 요청/수락/취소/완료", "완료 후 리뷰"]),
        ("Admin Dashboard", "검색, 필터, 테이블 중심의 운영 화면", "관리자 대시보드 / 상품 / 경매 / 출금 캡처 삽입", ["서버 관리자 권한 검증", "상태 배지와 빠른 액션", "모달/토스트 통일", "위험 액션 확인 모달"]),
    ]
    for heading, sub, ph, pts in screens:
        s = add(heading, "04 Screens")
        placeholder(s, 0.75, 1.45, 7.8, 4.75, ph)
        card(s, 8.9, 1.45, 3.55, 4.75, "핵심 포인트", "\n".join(pts), BLUE)

    # 14 Architecture
    s = add("System Architecture", "05 Architecture")
    flow(s, 0.85, 1.6, ["React", "REST API", "Spring Boot", "MySQL", "Redis"])
    card(s, 0.75, 2.65, 3.0, 2.4, "Frontend", "사용자 화면\n관리자 화면\n채팅 팝업", BLUE)
    card(s, 4.0, 2.65, 3.7, 2.4, "Backend", "Auth / Product / Auction / Trade\nChat / Mileage / Report / Admin\n권한 및 상태 검증", GREEN)
    card(s, 7.95, 2.65, 2.0, 2.4, "MySQL", "회원\n상품/경매\n거래/정산", ORANGE)
    card(s, 10.25, 2.65, 2.0, 2.4, "Redis", "이메일 인증\n캐시성 데이터", RED)
    card(s, 0.75, 5.45, 11.5, 0.8, "시스템 흐름", "React 요청 → Spring Boot 인증/인가 및 상태 검증 → MySQL 저장 → Redis 인증/캐시 → WebSocket 채팅", NAVY)

    # 15 API
    s = add("Core API", "06 API / DB")
    table(s, 0.75, 1.55, 11.8, 2.0, [["기능", "Method", "API", "설명"], ["이미지 업로드", "POST", "/api/products/images", "상품/경매 등록용 이미지 업로드 및 서버 검증"], ["경매 입찰", "POST", "/api/auctions/{auctionId}/bids", "경매 상태, 최고가, 마일리지 검증 후 입찰 처리"], ["출금 처리", "PATCH", "/api/admin/mileage/withdrawals/{id}", "관리자 출금 승인/반려 처리"]], [1.45, 0.95, 3.75, 5.65], 8.2)
    card(s, 0.75, 4.15, 11.8, 1.2, "API 설계 기준", "사용자 액션은 서버에서 권한과 상태를 재검증하고, 관리자 API는 프론트 메뉴 숨김과 별개로 서버에서 권한을 강제했습니다.", BLUE)

    # 16 DB
    s = add("DB Design", "06 API / DB")
    table(s, 0.75, 1.45, 6.35, 4.7, [["도메인", "주요 테이블", "설명"], ["Auth", "users", "회원 정보, 권한, 상태"], ["Product", "products, product_images", "상품, 이미지, 찜"], ["Auction", "auctions, bids", "경매, 입찰 내역"], ["Trade", "trades", "상품 거래 상태"], ["Chat", "chat_rooms, chat_messages", "채팅방, 메시지"], ["Mileage", "mileage_accounts, ledger, withdrawals", "잔액, 원장, 출금"], ["Admin", "reports, restrictions, audit_logs", "신고, 제재, 감사 로그"]], [1.15, 2.8, 2.4], 7.2)
    placeholder(s, 7.45, 1.45, 4.8, 3.25, "ERD 캡처 삽입")
    card(s, 7.45, 5.05, 4.8, 1.1, "설계 포인트", "상품/경매는 분리하고 사용자, 리뷰, 채팅 흐름은 연결했습니다.", GREEN)

    # 17 Section
    s = prs.slides.add_slide(blank); section(s, "07 LOGIC", "핵심 로직", "기능 설명보다 왜 필요했고 어떻게 해결했는지를 중심으로 정리했습니다."); page += 1

    # 18-21 Logic
    logic = [
        ("Auction Bid / Mileage", "경매 입찰 및 마일리지 예약", ["경매 입찰은 현재 최고가, 입찰자 권한, 경매 상태, 마일리지 잔액, 이전 입찰자의 예약 해제를 함께 처리해야 합니다.", "상태/마감/최소 입찰가/본인 상품 입찰/현재 최고 입찰자 여부를 검증했습니다.", "신규 최고 입찰자의 마일리지를 예약하고 이전 최고 입찰자의 예약 마일리지를 해제했습니다.", "최고가, 최고 입찰자, 예약 금액의 일관성을 유지했습니다."]),
        ("Access Control", "권한 기반 메뉴 및 API 접근 제어", ["프론트에서 버튼을 숨기는 것만으로는 보안이 보장되지 않습니다.", "사용자 역할에 따라 메뉴와 버튼을 제어하고, 관리자 API는 서버에서 접근을 강제했습니다.", "신고/출금/회원 제재 같은 관리자 액션은 서버에서 재검증했습니다.", "UI 편의성과 서버 보안을 분리하여 권한 없는 접근을 방지했습니다."]),
        ("Image Upload Validation", "이미지 업로드 서버 검증", ["프론트 검증만으로는 확장자를 위장한 파일 업로드를 막기 어렵습니다.", "개수, 용량, 확장자, Content-Type, 파일 시그니처를 서버에서 검증했습니다.", "JPEG/PNG/GIF/WEBP 시그니처를 확인하고 저장 경로도 검증했습니다.", "정상 업로드, 위장 파일 차단, 미지원 확장자 차단 테스트를 작성했습니다."]),
        ("React Data Loading", "React 데이터 로딩 흐름 안정화", ["Hook dependency warning은 데이터 로딩 구조를 점검하는 신호였습니다.", "조회 함수는 useCallback으로 안정화하고, 검색/필터는 URL query 기반으로 복원했습니다.", "경매 타이머는 auction 객체 전체가 아닌 id/status/endAt만 의존하도록 정리했습니다.", "warning 17개 제거 후 lint/build를 통과했습니다."]),
    ]
    for heading, sub, pts in logic:
        s = add(heading, "07 Logic")
        card(s, 0.75, 1.5, 11.8, 0.85, sub, "", BLUE)
        bullets(s, 1.0, 2.75, 11.0, 2.8, pts, 13)

    # 22-24 Troubleshooting
    trouble = [
        ("Troubleshooting 1", "경매 입찰 데이터 정합성", "여러 사용자가 입찰할 때 최고 입찰자, 현재 최고가, 마일리지 예약 상태가 불일치할 가능성이 있었습니다.", "입찰 등록, 최고가 갱신, 마일리지 예약/해제 처리가 서로 다른 도메인에 걸쳐 있었습니다.", "입찰 처리 흐름을 하나의 도메인 흐름으로 묶고, 경매 상태와 최고가를 재검증한 뒤 예약/해제를 처리했습니다.", "동시 입찰 상황에서도 최고 입찰자, 현재 최고가, 예약 금액이 일관되게 유지되도록 개선했습니다."),
        ("Troubleshooting 2", "관리자 권한 서버 강제", "관리자 메뉴와 버튼을 숨기더라도 URL 직접 접근이나 API 직접 호출로 관리자 기능 접근을 시도할 수 있었습니다.", "프론트 UI 제어는 사용자 경험 측면의 처리일 뿐 실제 보안 경계는 서버에서 강제되어야 했습니다.", "Spring Security에서 관리자 API 경로를 제한하고, 프론트에서는 권한별 메뉴/버튼 노출을 분리했습니다.", "권한 없는 사용자의 직접 접근과 API 호출이 서버에서 차단되도록 개선했습니다."),
        ("Troubleshooting 3", "경매 마감 상태 불일치", "경매 종료 시간이 지난 직후 서버 상태 갱신 전에는 화면상 입찰 가능 상태처럼 보일 수 있었습니다.", "경매 상태는 서버의 OPEN/CLOSED 값과 클라이언트 현재 시간 기준 마감 여부를 함께 고려해야 했습니다.", "프론트에서는 endAt 기준으로 버튼과 상태 배지를 즉시 비활성화하고, 서버 입찰 API에서도 상태와 마감 시간을 재검증했습니다.", "마감 직후에도 화면 상태와 실제 입찰 가능 여부가 일관되도록 개선했습니다."),
    ]
    for heading, sub, problem, cause, solution, result in trouble:
        s = add(heading, "08 Trouble")
        items = [("문제", problem, RED), ("원인", cause, ORANGE), ("해결", solution, BLUE), ("결과", result, GREEN)]
        y = 1.45
        for label, body, color in items:
            card(s, 0.85, y, 11.6, 0.9, label, body, color)
            y += 1.18

    # 25 Test
    s = add("Test / Verification", "09 Verify")
    table(s, 0.75, 1.55, 11.8, 2.3, [["구분", "검증"], ["Frontend", "npm run lint, npm run build"], ["Backend", "이미지 업로드 통합 테스트"], ["Image Upload", "정상 PNG 업로드, 위장 파일 차단, SVG 차단"], ["Hook Cleanup", "기존 17개 warning 제거 후 빌드 성공"]], [2.1, 9.7], 9)
    card(s, 0.75, 4.4, 11.8, 1.05, "남은 보완", "전체 백엔드 테스트는 Redis 연결 환경이 필요하며, 결제/출금/경매 동시성 관련 테스트를 추가할 수 있습니다.", ORANGE)

    # 26 Retrospective
    s = add("Retrospective", "10 Review")
    card(s, 0.75, 1.45, 3.7, 4.5, "배운 점", "연결된 기능은 상태 전이 설계가 중요\n권한 제어는 서버에서 강제해야 함\nHook warning은 데이터 로딩 구조 점검 신호", BLUE)
    card(s, 4.85, 1.45, 3.7, 4.5, "아쉬운 점", "전체 E2E 테스트 부족\nRedis 외부 의존성 테스트 환경 분리 필요\n배포 환경 모니터링/로그 설계 보완 필요", ORANGE)
    card(s, 8.95, 1.45, 3.4, 4.5, "개선 방향", "경매 동시성 테스트 강화\n테스트 프로파일 분리\n관리자 통계 고도화\nE2E 테스트 추가", GREEN)

    # 27 Future
    s = add("Future Improvements", "11 Future")
    bullets(s, 1.0, 1.65, 11.0, 3.9, ["배포 환경 구성 및 CI/CD 적용", "경매 입찰 동시성 테스트 강화", "관리자 통계 시각화 고도화", "검색 성능 개선 및 인덱싱 전략 보완", "알림 기능 실시간화", "주요 사용자 흐름 E2E 테스트 추가"], 15)

    # 28 Closing
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    text(s, 0.95, 1.0, 9.5, 0.6, "Thank You", 34, True, WHITE)
    text(s, 0.98, 1.85, 10.8, 0.5, "단순 기능 구현을 넘어 사용자 거래 흐름과 관리자 운영 흐름을 연결한 프로젝트입니다.", 15, False, RGBColor(203, 213, 225))
    table(s, 1.0, 3.0, 7.0, 2.0, [["항목", "내용"], ["GitHub", ""], ["Blog / Notion", ""], ["시연 영상", ""], ["연락처", ""]], [1.6, 5.4], 10.5)
    return prs


if __name__ == "__main__":
    prs_new().save(OUT)
    print(OUT)
