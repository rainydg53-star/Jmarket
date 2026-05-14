from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = "Jmarket_Portfolio_v3.pptx"
FONT = "Malgun Gothic"

INK = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
SUBTLE = RGBColor(100, 116, 139)
LINE = RGBColor(203, 213, 225)
PANEL = RGBColor(248, 250, 252)
PANEL_2 = RGBColor(241, 245, 249)
NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
BLUE_LIGHT = RGBColor(219, 234, 254)
GREEN = RGBColor(5, 150, 105)
GREEN_LIGHT = RGBColor(209, 250, 229)
ORANGE = RGBColor(234, 88, 12)
ORANGE_LIGHT = RGBColor(255, 237, 213)
RED = RGBColor(220, 38, 38)
RED_LIGHT = RGBColor(254, 226, 226)
WHITE = RGBColor(255, 255, 255)


def fmt(run, size=11, bold=False, color=INK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def tx(slide, x, y, w, h, value, size=11, bold=False, color=INK, align=None, valign=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    if valign:
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = value
    fmt(r, size, bold, color)
    return box


def shape(slide, x, y, w, h, fill=WHITE, border=LINE, radius=True):
    s = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(0.7)
    return s


def topbar(slide, page, section, title, subtitle=None):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.82))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    tx(slide, 0.55, 0.16, 1.55, 0.24, section, 8.5, True, RGBColor(147, 197, 253))
    tx(slide, 0.55, 0.88, 8.5, 0.42, title, 21, True, NAVY)
    if subtitle:
        tx(slide, 0.58, 1.31, 8.8, 0.22, subtitle, 8.8, False, SUBTLE)
    tx(slide, 11.8, 0.2, 0.8, 0.24, f"{page:02}", 8.5, True, WHITE, PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.62), Inches(12.2), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def footer(slide):
    tx(slide, 0.55, 7.08, 3.2, 0.18, "Jmarket Portfolio", 7.5, False, SUBTLE, valign=False)


def bullet_list(slide, x, y, w, h, items, size=10.5, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.space_after = Pt(4)
        p._p.get_or_add_pPr().set("marL", "190500")
        p._p.get_or_add_pPr().set("indent", "-142875")
        for r in p.runs:
            fmt(r, size, False, color)
    return box


def card(slide, x, y, w, h, title, body=None, accent=BLUE, fill=WHITE):
    shape(slide, x, y, w, h, fill, LINE)
    head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.42))
    head.fill.solid()
    head.fill.fore_color.rgb = accent
    head.line.fill.background()
    tx(slide, x + 0.18, y + 0.08, w - 0.36, 0.2, title, 9.5, True, WHITE, valign=False)
    if body:
        tx(slide, x + 0.22, y + 0.56, w - 0.44, h - 0.68, body, 10.2, False, INK, valign=False)


def placeholder(slide, x, y, w, h, label):
    shape(slide, x, y, w, h, PANEL_2, LINE)
    tx(slide, x + 0.25, y + h / 2 - 0.18, w - 0.5, 0.36, label, 13, True, SUBTLE, PP_ALIGN.CENTER)


def table(slide, x, y, w, h, rows, widths=None, size=8.5):
    tbl_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = tbl_shape.table
    if widths:
        for idx, width in enumerate(widths):
            tbl.columns[idx].width = Inches(width)
    row_h = Inches(h / len(rows))
    for idx in range(len(rows)):
        tbl.rows[idx].height = row_h
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = value
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.fore_color.rgb = RGBColor(248, 250, 252) if r_idx % 2 else RGBColor(226, 232, 240)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    fmt(run, size, r_idx == 0, WHITE if r_idx == 0 else INK)
    return tbl_shape


def metric(slide, x, y, w, label, value, accent=BLUE):
    shape(slide, x, y, w, 0.86, WHITE, LINE)
    tx(slide, x + 0.18, y + 0.13, w - 0.36, 0.18, label, 8, True, SUBTLE, valign=False)
    tx(slide, x + 0.18, y + 0.42, w - 0.36, 0.28, value, 16, True, accent, valign=False)


def process(slide, x, y, items, accent=BLUE):
    step_w = 1.62
    for idx, item in enumerate(items):
        shape(slide, x + idx * 2.05, y, step_w, 0.62, WHITE, LINE)
        tx(slide, x + idx * 2.05 + 0.1, y + 0.17, step_w - 0.2, 0.2, item, 8.8, True, NAVY, PP_ALIGN.CENTER)
        if idx < len(items) - 1:
            tx(slide, x + idx * 2.05 + step_w + 0.06, y + 0.18, 0.25, 0.18, "→", 11, True, accent)


def section_slide(prs, label, title, body):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    tx(slide, 0.95, 1.25, 2.3, 0.25, label, 9, True, RGBColor(147, 197, 253))
    tx(slide, 0.95, 1.78, 8.8, 0.7, title, 30, True, WHITE)
    tx(slide, 0.98, 2.75, 8.5, 0.45, body, 13.5, False, RGBColor(203, 213, 225))
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.95), Inches(3.65), Inches(1.35), Inches(0.055))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    page = 1

    def slide(section, title, subtitle=None):
        nonlocal page
        s = prs.slides.add_slide(blank)
        topbar(s, page, section, title, subtitle)
        footer(s)
        page += 1
        return s

    # 1
    s = prs.slides.add_slide(blank)
    left = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(3.75), Inches(7.5))
    left.fill.solid(); left.fill.fore_color.rgb = NAVY; left.line.fill.background()
    tx(s, 0.62, 0.78, 2.5, 0.24, "PORTFOLIO", 9, True, RGBColor(147, 197, 253))
    tx(s, 0.62, 1.28, 2.7, 0.52, "Jmarket", 30, True, WHITE)
    tx(s, 0.64, 1.95, 2.55, 0.62, "중고거래 및 경매 기반 마켓 플랫폼", 12.5, False, RGBColor(203, 213, 225), valign=False)
    tx(s, 4.55, 1.1, 7.5, 0.7, "사용자 거래 흐름과 관리자 운영 흐름을 연결한 개인 프로젝트", 25, True, NAVY, valign=False)
    tx(s, 4.58, 2.0, 7.8, 0.32, "Spring Boot · React · WebSocket · MySQL · Redis", 13, False, MUTED)
    metric(s, 4.6, 3.0, 2.25, "Role", "Full-stack")
    metric(s, 7.1, 3.0, 2.25, "Focus", "정합성", GREEN)
    metric(s, 9.6, 3.0, 2.25, "Domain", "거래/경매", ORANGE)
    tx(s, 4.6, 6.32, 5.5, 0.25, "이름 / 연락처 / 이메일", 10, False, SUBTLE)
    page += 1

    # 2
    s = slide("01 PROFILE", "Profile", "인적사항과 간단 소개")
    table(s, 0.75, 1.95, 5.4, 2.65, [["항목", "내용"], ["이름", ""], ["연락처", ""], ["이메일", ""], ["GitHub", ""], ["Blog / Notion", ""]], [1.25, 4.15], 8.8)
    card(s, 6.55, 1.95, 5.95, 2.65, "한 줄 소개", "Spring Boot와 React 기반으로 거래, 경매, 채팅, 마일리지, 관리자 운영 기능을 구현했습니다. 기능 구현뿐 아니라 권한, 상태 전이, 데이터 정합성을 함께 고려했습니다.", BLUE)
    table(s, 0.75, 5.05, 5.4, 1.25, [["기간", "기관", "내용"], ["", "", ""], ["", "", ""]], [1.25, 1.7, 2.45], 8.4)
    table(s, 6.55, 5.05, 5.95, 1.25, [["취득일", "자격증", "기관"], ["", "", ""], ["", "", ""]], [1.35, 2.5, 2.1], 8.4)

    # 3
    s = slide("01 PROFILE", "Education / Skills", "기술 스택은 사용 이유가 드러나도록 정리")
    table(s, 0.75, 1.95, 5.55, 1.35, [["기간", "기관", "내용"], ["", "", ""], ["", "", ""]], [1.15, 1.6, 2.8], 8.3)
    table(s, 0.75, 3.75, 5.55, 1.25, [["취득일", "자격증", "기관"], ["", "", ""], ["", "", ""]], [1.25, 2.4, 1.9], 8.3)
    table(s, 6.65, 1.95, 5.85, 3.05, [
        ["구분", "기술"],
        ["Backend", "Java 21, Spring Boot, Spring Security, JPA"],
        ["Frontend", "React, Vite, JavaScript, CSS"],
        ["DB / Cache", "MySQL, Redis"],
        ["Realtime", "WebSocket, STOMP"],
        ["Test / Tool", "JUnit, MockMvc, Gradle, npm, Git"],
    ], [1.45, 4.4], 7.6)
    card(s, 0.75, 5.45, 11.75, 0.95, "개인 학습 / 활동", "인증/인가, 경매 정산, React Hook 의존성, 관리자 운영 화면 중심으로 학습 및 구현", GREEN, PANEL)

    # 4
    s = slide("02 CONTENTS", "Project Contents", "포트폴리오에 포함할 프로젝트")
    table(s, 0.75, 1.95, 11.75, 1.35, [["프로젝트명", "기간", "인원", "기술", "설명"], ["Jmarket", "2026.04 ~ 2026.05", "개인", "Spring Boot, React, MySQL, Redis, WebSocket", "중고거래와 경매를 함께 제공하는 마켓 플랫폼"]], [1.45, 1.6, 0.8, 3.55, 4.35], 8.2)
    placeholder(s, 0.75, 3.75, 11.75, 2.25, "프로젝트 대표 화면 또는 서비스 흐름 이미지 삽입")

    section_slide(prs, "03 PROJECT", "Jmarket 상세", "문제 해결 중심으로 프로젝트의 구조와 핵심 로직을 정리했습니다.")
    page += 1

    # 6 overview
    s = slide("03 PROJECT", "Project Overview", "OO 문제를 해결하기 위해 OO 시스템 개발")
    card(s, 0.75, 1.85, 11.75, 0.95, "핵심 문장", "중고거래 과정에서 발생하는 상품 탐색, 거래 요청, 경매 입찰, 채팅, 마일리지 정산, 관리자 운영 문제를 해결하기 위해 통합 마켓 플랫폼을 개발했습니다.", BLUE, PANEL)
    card(s, 0.75, 3.25, 3.6, 2.4, "해결하려는 문제", "거래 요청부터 완료까지 상태 관리\n경매 입찰과 마일리지 예약 정합성\n신고/출금/회원 제재 운영 분리", RED)
    card(s, 4.85, 3.25, 3.6, 2.4, "기대 효과", "사용자는 한 흐름에서 거래 처리\n관리자는 운영 데이터를 빠르게 확인\n권한/상태별 액션으로 오류 감소", GREEN)
    card(s, 8.95, 3.25, 3.55, 2.4, "차별화 포인트", "일반 거래와 경매 동시 제공\n마일리지 예약/해제/정산 포함\n사용자/관리자 UX 분리", ORANGE)

    # 7 main features
    s = slide("03 PROJECT", "Main Features", "핵심 기능 요약")
    features = ["회원가입, 로그인, 이메일 인증, 비밀번호 찾기/변경", "상품 등록/수정/삭제, 이미지 드래그 업로드, 검색/필터", "경매 등록, 입찰, 즉시구매, 입찰 내역, 마감 처리", "거래 요청/수락/취소/완료 및 리뷰", "실시간 채팅방 목록/팝업", "마일리지 충전/사용/출금 요청 및 관리자 처리", "신고, 고객센터 문의, 알림", "관리자 대시보드와 운영 관리"]
    for i, feat in enumerate(features):
        x = 0.85 + (i % 2) * 5.95
        y = 1.95 + (i // 2) * 1.05
        shape(s, x, y, 5.55, 0.72, WHITE, LINE)
        tx(s, x + 0.18, y + 0.2, 0.35, 0.2, f"{i+1}", 10, True, BLUE)
        tx(s, x + 0.58, y + 0.17, 4.75, 0.22, feat, 10.2, False, INK)

    # 8 role
    s = slide("03 PROJECT", "Role & Responsibility", "개인 프로젝트 담당 범위")
    metric(s, 0.75, 1.85, 2.6, "Project", "개인", BLUE)
    metric(s, 3.6, 1.85, 2.6, "Hook Warning", "17 → 0", GREEN)
    metric(s, 6.45, 1.85, 2.6, "Core Domain", "8+", ORANGE)
    metric(s, 9.3, 1.85, 2.6, "Slides", "28p", BLUE)
    bullet_list(s, 0.95, 3.15, 11.1, 2.5, ["기획, DB 설계, 백엔드 API, 프론트 화면 구현 전담", "Spring Security 기반 인증/인가와 관리자 API 접근 제한", "상품/경매/거래/채팅/마일리지 도메인 구현", "WebSocket/STOMP 채팅과 관리자 운영 화면 개선", "이미지 업로드 서버 검증 및 통합 테스트 작성"], 12)

    # 9-13 screen
    screens = [
        ("04 SCREENS", "Main Page", "비로그인 사용자도 접근 가능한 메인 화면", "메인 화면 캡처 삽입", ["실시간 검색어 순위", "급상승 물품", "빈 상태 UI", "상품/경매 탐색 흐름"]),
        ("04 SCREENS", "Product Flow", "상품 목록/상세/등록 흐름", "상품 목록 / 상세 / 등록 화면 캡처 삽입", ["검색/필터 URL 유지", "썸네일 클릭 상세 이동", "드래그 이미지 업로드", "서버 이미지 검증"]),
        ("04 SCREENS", "Auction Flow", "경매 목록/상세/입찰 흐름", "경매 목록 / 상세 / 입찰 내역 화면 캡처 삽입", ["상태별 입찰 가능 여부 제어", "최고 입찰자와 최고가 표시", "입찰 내역 차트", "마감 타이머 상태 갱신"]),
        ("04 SCREENS", "Chat / Trade Flow", "채팅과 거래 상태 흐름", "채팅 목록 / 팝업 / 거래 목록 화면 캡처 삽입", ["WebSocket/STOMP 실시간 메시지", "채팅방 새 창 팝업", "거래 요청/수락/취소/완료", "거래 완료 후 리뷰 작성"]),
        ("04 SCREENS", "Admin Dashboard", "운영 도구형 관리자 화면", "관리자 대시보드 / 상품 / 경매 / 출금 화면 캡처 삽입", ["서버 관리자 권한 검증", "검색/필터/테이블/상태 배지", "액션 결과 모달/토스트", "위험 액션 확인 모달"]),
    ]
    for sec, ttl, sub, ph, pts in screens:
        s = slide(sec, ttl, sub)
        placeholder(s, 0.75, 1.95, 7.65, 4.45, ph)
        card(s, 8.8, 1.95, 3.7, 4.45, "핵심 포인트", "\n".join(pts), BLUE)

    # 14 architecture
    s = slide("05 ARCHITECTURE", "System Architecture", "Frontend / Backend / DB / Infra 구조")
    process(s, 0.85, 2.0, ["React", "REST API", "Spring Boot", "MySQL", "Redis"])
    card(s, 0.75, 3.0, 2.75, 2.25, "Frontend", "사용자 화면\n관리자 화면\n채팅 팝업", BLUE)
    card(s, 3.85, 3.0, 3.65, 2.25, "Backend", "Auth / Product / Auction / Trade\nChat / Mileage / Report / Admin\n권한 및 상태 검증", GREEN)
    card(s, 7.85, 3.0, 2.0, 2.25, "MySQL", "회원\n상품/경매\n거래/정산", ORANGE)
    card(s, 10.2, 3.0, 2.0, 2.25, "Redis", "이메일 인증\n캐시성 데이터", RED)
    card(s, 0.75, 5.75, 11.45, 0.78, "시스템 흐름", "React 요청 → Spring Boot 인증/인가 및 상태 검증 → MySQL 저장 → Redis 인증/캐시 → WebSocket 채팅", NAVY)

    # 15 API
    s = slide("06 API / DB", "Core API", "전체 API가 아닌 핵심 흐름 중심")
    table(s, 0.75, 1.95, 11.75, 2.0, [["기능", "Method", "API", "설명"], ["이미지 업로드", "POST", "/api/products/images", "상품/경매 등록용 이미지 업로드 및 서버 검증"], ["경매 입찰", "POST", "/api/auctions/{auctionId}/bids", "경매 상태, 최고가, 마일리지 검증 후 입찰 처리"], ["출금 처리", "PATCH", "/api/admin/mileage/withdrawals/{id}", "관리자 출금 승인/반려 처리"]], [1.45, 0.95, 3.75, 5.6], 8.0)
    card(s, 0.75, 4.45, 11.75, 1.05, "API 설계 기준", "사용자 액션은 서버에서 권한과 상태를 재검증하고, 관리자 API는 프론트 메뉴 숨김과 별개로 서버에서 권한을 강제했습니다.", BLUE, PANEL)

    # 16 DB
    s = slide("06 API / DB", "DB Design", "주요 도메인과 테이블 구조")
    table(s, 0.75, 1.95, 6.35, 4.4, [["도메인", "주요 테이블", "설명"], ["Auth", "users", "회원 정보, 권한, 상태"], ["Product", "products, product_images", "상품, 이미지, 찜"], ["Auction", "auctions, bids", "경매, 입찰 내역"], ["Trade", "trades", "상품 거래 상태"], ["Chat", "chat_rooms, chat_messages", "채팅방, 메시지"], ["Mileage", "mileage_accounts, ledger, withdrawals", "잔액, 원장, 출금"], ["Admin", "reports, restrictions, audit_logs", "신고, 제재, 감사 로그"]], [1.15, 2.75, 2.45], 7.0)
    placeholder(s, 7.45, 1.95, 4.75, 3.0, "ERD 캡처 삽입")
    card(s, 7.45, 5.35, 4.75, 0.95, "설계 포인트", "상품/경매는 분리하고 사용자, 리뷰, 채팅 흐름은 연결했습니다.", GREEN, PANEL)

    section_slide(prs, "07 LOGIC", "핵심 로직", "기능 설명보다 왜 필요했고 어떻게 해결했는지를 중심으로 정리했습니다.")
    page += 1

    # 18-21 logic
    logic = [
        ("07 LOGIC", "Auction Bid / Mileage", "경매 입찰 및 마일리지 예약", ["경매 입찰은 현재 최고가, 입찰자 권한, 경매 상태, 마일리지 잔액, 이전 입찰자의 예약 해제를 함께 처리해야 합니다.", "상태/마감/최소 입찰가/본인 상품 입찰/현재 최고 입찰자 여부를 검증했습니다.", "신규 최고 입찰자의 마일리지를 예약하고 이전 최고 입찰자의 예약 마일리지를 해제했습니다.", "최고가, 최고 입찰자, 예약 금액의 일관성을 유지했습니다."]),
        ("07 LOGIC", "Access Control", "권한 기반 메뉴 및 API 접근 제어", ["프론트에서 버튼을 숨기는 것만으로는 보안이 보장되지 않습니다.", "사용자 역할에 따라 메뉴와 버튼을 제어하고, 관리자 API는 서버에서 접근을 강제했습니다.", "신고/출금/회원 제재 같은 관리자 액션은 서버에서 재검증했습니다.", "UI 편의성과 서버 보안을 분리하여 권한 없는 접근을 방지했습니다."]),
        ("07 LOGIC", "Image Upload Validation", "이미지 업로드 서버 검증", ["프론트 검증만으로는 확장자를 위장한 파일 업로드를 막기 어렵습니다.", "개수, 용량, 확장자, Content-Type, 파일 시그니처를 서버에서 검증했습니다.", "JPEG/PNG/GIF/WEBP 시그니처를 확인하고 저장 경로도 검증했습니다.", "정상 업로드, 위장 파일 차단, 미지원 확장자 차단 테스트를 작성했습니다."]),
        ("07 LOGIC", "React Data Loading", "React 데이터 로딩 흐름 안정화", ["Hook dependency warning은 데이터 로딩 구조를 점검하는 신호였습니다.", "조회 함수는 useCallback으로 안정화하고, 검색/필터는 URL query 기반으로 복원했습니다.", "경매 타이머는 auction 객체 전체가 아닌 id/status/endAt만 의존하도록 정리했습니다.", "warning 17개 제거 후 lint/build를 통과했습니다."]),
    ]
    for sec, ttl, heading, pts in logic:
        s = slide(sec, ttl, heading)
        card(s, 0.75, 1.95, 11.75, 4.45, heading, None, BLUE, PANEL)
        bullet_list(s, 1.08, 2.75, 10.95, 2.85, pts, 12.4)

    # 22-24 troubleshooting
    trouble = [
        ("08 TROUBLE", "Troubleshooting 1", "경매 입찰 데이터 정합성", "여러 사용자가 입찰할 때 최고 입찰자, 현재 최고가, 마일리지 예약 상태가 불일치할 가능성이 있었습니다.", "입찰 등록, 최고가 갱신, 마일리지 예약/해제 처리가 서로 다른 도메인에 걸쳐 있었습니다.", "입찰 처리 흐름을 하나의 도메인 흐름으로 묶고, 경매 상태와 최고가를 재검증한 뒤 예약/해제를 처리했습니다.", "동시 입찰 상황에서도 최고 입찰자, 현재 최고가, 예약 금액이 일관되게 유지되도록 개선했습니다."),
        ("08 TROUBLE", "Troubleshooting 2", "관리자 권한 서버 강제", "관리자 메뉴와 버튼을 숨기더라도 URL 직접 접근이나 API 직접 호출로 관리자 기능 접근을 시도할 수 있었습니다.", "프론트 UI 제어는 사용자 경험 측면의 처리일 뿐 실제 보안 경계는 서버에서 강제되어야 했습니다.", "Spring Security에서 관리자 API 경로를 제한하고, 프론트에서는 권한별 메뉴/버튼 노출을 분리했습니다.", "권한 없는 사용자의 직접 접근과 API 호출이 서버에서 차단되도록 개선했습니다."),
        ("08 TROUBLE", "Troubleshooting 3", "경매 마감 상태 불일치", "경매 종료 시간이 지난 직후 서버 상태 갱신 전에는 화면상 입찰 가능 상태처럼 보일 수 있었습니다.", "경매 상태는 서버의 OPEN/CLOSED 값과 클라이언트 현재 시간 기준 마감 여부를 함께 고려해야 했습니다.", "프론트에서는 endAt 기준으로 버튼과 상태 배지를 즉시 비활성화하고, 서버 입찰 API에서도 상태와 마감 시간을 재검증했습니다.", "마감 직후에도 화면 상태와 실제 입찰 가능 여부가 일관되도록 개선했습니다."),
    ]
    for sec, ttl, sub, problem, cause, solution, result in trouble:
        s = slide(sec, ttl, sub)
        card(s, 0.75, 1.9, 11.75, 0.9, "문제 상황", problem, RED)
        card(s, 0.75, 3.05, 11.75, 0.9, "원인", cause, ORANGE)
        card(s, 0.75, 4.2, 11.75, 0.9, "해결 방법", solution, BLUE)
        card(s, 0.75, 5.35, 11.75, 0.9, "결과", result, GREEN)

    # 25 verification
    s = slide("09 VERIFY", "Test / Verification", "구현 후 검증한 내용")
    table(s, 0.75, 1.95, 11.75, 2.35, [["구분", "검증"], ["Frontend", "npm run lint, npm run build"], ["Backend", "이미지 업로드 통합 테스트"], ["Image Upload", "정상 PNG 업로드, 위장 파일 차단, SVG 차단"], ["Hook Cleanup", "기존 17개 warning 제거 후 빌드 성공"]], [2.05, 9.7], 8.8)
    card(s, 0.75, 4.85, 11.75, 0.95, "남은 보완", "전체 백엔드 테스트는 Redis 연결 환경이 필요하며, 결제/출금/경매 동시성 관련 테스트를 추가할 수 있습니다.", ORANGE, PANEL)

    # 26 review
    s = slide("10 REVIEW", "Retrospective", "배운 점, 아쉬운 점, 개선 방향")
    card(s, 0.75, 1.95, 3.65, 3.95, "배운 점", "연결된 기능은 상태 전이 설계가 중요\n권한 제어는 서버에서 강제해야 함\nHook warning은 데이터 로딩 구조 점검 신호", BLUE)
    card(s, 4.85, 1.95, 3.65, 3.95, "아쉬운 점", "전체 E2E 테스트 부족\nRedis 외부 의존성 테스트 환경 분리 필요\n배포 환경 모니터링/로그 설계 보완 필요", ORANGE)
    card(s, 8.95, 1.95, 3.35, 3.95, "개선 방향", "경매 동시성 테스트 강화\n테스트 프로파일 분리\n관리자 통계 고도화\nE2E 테스트 추가", GREEN)

    # 27 future
    s = slide("11 FUTURE", "Future Improvements", "추가 개선 방향")
    bullet_list(s, 1.0, 1.95, 11.0, 3.9, ["배포 환경 구성 및 CI/CD 적용", "경매 입찰 동시성 테스트 강화", "관리자 통계 시각화 고도화", "검색 성능 개선 및 인덱싱 전략 보완", "알림 기능 실시간화", "주요 사용자 흐름 E2E 테스트 추가"], 14)

    # 28 closing
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    tx(s, 0.95, 1.05, 8.0, 0.6, "Thank You", 34, True, WHITE, valign=False)
    tx(s, 0.98, 1.88, 10.8, 0.38, "단순 기능 구현을 넘어 사용자 거래 흐름과 관리자 운영 흐름을 연결한 프로젝트입니다.", 14.5, False, RGBColor(203, 213, 225))
    table(s, 1.0, 3.05, 7.0, 2.0, [["항목", "내용"], ["GitHub", ""], ["Blog / Notion", ""], ["시연 영상", ""], ["연락처", ""]], [1.55, 5.45], 10.2)

    return prs


if __name__ == "__main__":
    build().save(OUT)
    print(OUT)
